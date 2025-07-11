from pptx import Presentation

def pptx_to_txt(pptx_path, txt_path):
    # Carrega a apresentação
    prs = Presentation(pptx_path)
    texto_completo = []

    # Itera por cada slide
    for slide in prs.slides:
        for shape in slide.shapes:
            # Extrai texto de formas com texto
            if hasattr(shape, "text") and shape.text.strip():
                texto_completo.append(shape.text.strip())
            # Extrai texto de tabelas
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            texto_completo.append(cell.text.strip())

    # Salva o texto em um arquivo .txt
    with open(txt_path, 'w', encoding='utf-8') as f:
        for linha in texto_completo:
            f.write(linha + '\n')

# Converte o arquivo
pptx_to_txt('emis-ptbr.pptx', 'emis-ptbr.txt')